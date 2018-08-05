from pptx import Presentation

# CONSTANTS
# HEIGHT = 9
# WIDTH = 16
# LEFTMOST = Inches(0)
# TOPMOST = Inches(0)
# HEIGHT_IN = Inches(HEIGHT)
# WIDTH_IN = Inches(WIDTH)

# One inch equates to 914400 EMUs
# INCHES_TO_EMU = 914400
# One centimeter is 360000 EMUs
# CMS_TO_EMU = 360000

# Location of powerpoint template
POWERPOINT_TEMPLATE_FILE = 'data/powerpoint/template.pptx'

# Layouts index in template
LAYOUT_TITLE_SLIDE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_SECTION_HEADER = 2
LAYOUT_TWO_CONTENT = 3
LAYOUT_TWO_TITLE_AND_CONTENT = 4
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6
LAYOUT_CONTENT_CAPTION = 7
LAYOUT_PICTURE_CAPTION = 8
LAYOUT_FULL_PICTURE = 11
LAYOUT_TITLE_AND_PICTURE = 12
LAYOUT_LARGE_QUOTE = 13
LAYOUT_TWO_TITLE_AND_IMAGE = 14


# HELPERS
def _create_slide(prs, slide_type):
    """ Creates a new slide in the given presentation using the slide_type template """
    return prs.slides.add_slide(prs.slide_layouts[slide_type])


def _add_title(slide, title):
    """ Adds the given title to the slide if the title is present"""
    if title:
        title_object = slide.shapes.title
        title_object.text = title


def _add_image(slide, placeholder_id, image_url, fit_image=True):
    placeholder = slide.placeholders[placeholder_id]
    if fit_image:
        pic = slide.shapes.add_picture(image_url, placeholder.left, placeholder.top)

        # calculate max width/height for target size
        ratio = min(placeholder.width / float(pic.width), placeholder.height / float(pic.height))

        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)

        pic.left = int(placeholder.left + ((placeholder.width - pic.width) / 2))
        pic.top = int(placeholder.top + ((placeholder.height - pic.height) / 2))

        placeholder = placeholder.element
        placeholder.getparent().remove(placeholder)
    else:
        placeholder.insert_picture(image_url)


def _add_text(slide, placeholder_id, text):
    placeholder = slide.placeholders[placeholder_id]
    placeholder.text = text


# FORMAT GENERATORS
# These are functions that get some inputs (texts, images...)
# and create layouted slides with these inputs

def create_new_powerpoint():
    return Presentation(POWERPOINT_TEMPLATE_FILE)


def create_title_slide(prs, title):
    slide = _create_slide(prs, LAYOUT_TITLE_SLIDE)
    _add_title(slide, title)
    return slide


def create_large_quote_slide(prs, text):
    if bool(text):
        slide = _create_slide(prs, LAYOUT_LARGE_QUOTE)
        _add_text(slide, 1, text)
        return slide


def create_image_slide(prs, title=None, image_url=None):
    """ Creates a slide with an image covering the whole slide"""
    # TODO debug this: the image can not be set!
    return _create_single_image_slide(prs, title, image_url, LAYOUT_TITLE_AND_PICTURE, True)


def create_full_image_slide(prs, title=None, image_url=None):
    """ Creates a slide with an image covering the whole slide"""
    return _create_single_image_slide(prs, title, image_url, LAYOUT_FULL_PICTURE, False)


def create_two_column_images_slide(prs, title=None, caption_1=None, image_1=None, caption_2=None, image_2=None):
    if bool(image_1) and bool(image_2):
        slide = _create_slide(prs, LAYOUT_TWO_TITLE_AND_IMAGE)
        _add_title(slide, title)
        _add_text(slide, 1, caption_1)
        _add_image(slide, 13, image_1, False)
        _add_text(slide, 3, caption_2)
        _add_image(slide, 14, image_2, False)
        return slide


def _create_single_image_slide(prs, title, image_url, slide_template_idx, fit_image):
    if image_url:
        slide = _create_slide(prs, slide_template_idx)
        _add_title(slide, title)
        _add_image(slide, 1, image_url, fit_image)
        return slide


# GENERATORS: Same as the template fillers above, but using generation functions

def generate_full_image_slide(title_generator, image_generator):
    return lambda prs, seed: create_full_image_slide(prs, title_generator(seed), image_generator(seed))


def generate_image_slide(title_generator, image_generator):
    return lambda prs, seed: create_image_slide(prs, title_generator(seed), image_generator(seed))


def generate_title_slide(title_generator):
    return lambda prs, seed: create_title_slide(prs, title_generator(seed))


def generate_large_quote_slide(text_generator):
    return lambda prs, seed: create_large_quote_slide(prs, text_generator(seed))


def generate_two_column_images_slide(title_generator, caption_1_generator, image_1_generator, caption_2_generator,
                                     image_2_generator):
    return lambda prs, seed: create_two_column_images_slide(prs, title_generator(seed), caption_1_generator(seed),
                                                            image_1_generator(seed), caption_2_generator(seed),
                                                            image_2_generator(seed))


class TupledGenerator:
    """ This class got introduced due to Python not having multi-line lambdas,
    thus requiring us in order to use tupled generators"""

    def __init__(self, generator, rest_of_function):
        self._generator = generator
        self._rest_of_function = rest_of_function

    def generate(self, prs, seed):
        generated_tuple = self._generator(seed)
        return self._rest_of_function(prs, seed, generated_tuple)


def generate_two_column_images_slide_tuple_caption(title_generator, captions_generator, image_1_generator,
                                                   image_2_generator):
    tupled_generator = TupledGenerator(captions_generator,
                                       lambda prs, seed, tuple: create_two_column_images_slide(prs, title_generator(
                                           seed), tuple[0], image_1_generator(
                                           seed), tuple[1], image_2_generator(seed)))
    return tupled_generator.generate
